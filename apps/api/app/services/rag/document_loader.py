"""
document_loader.py — Multi-Format Text Extraction & Text Chunking
=================================================================
WHAT THIS DOES:
  1. Multi-Format Text Extraction: Converts binary document files (PDF, DOCX, PPTX, CSV, TXT, MD)
     into plain structured text strings.
  2. Sliding Window Chunking: Divides large text documents into overlapping segments
     optimized for embedding generation and retrieval in RAG pipelines.

WHY CHUNKING & OVERLAP ARE CRUCIAL (INTERVIEW TOPIC):
  - LLM Context Constraints: Embedding models & LLMs have finite context limits (e.g. 512, 8k tokens).
  - Retrieval Granularity: Smaller chunks allow precise vector similarity matching without pulling
    irrelevant noise from a 100-page document.
  - Chunk Overlap: Guarantees that sentences or thoughts sitting on chunk boundaries are preserved
    in both adjacent chunks rather than being cut mid-thought.
"""

import io
from pathlib import Path
import pandas as pd
from fastapi import HTTPException, status

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".txt", ".md", ".pptx", ".csv"}


def extract_text_from_bytes(filename: str, data: bytes) -> str:
    """
    Extract readable text content from raw bytes according to file extension.
    """
    suffix = Path(filename).suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported file format '{suffix}'. Supported formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )

    if not data:
        return ""

    if suffix in {".txt", ".md"}:
        try:
            return data.decode("utf-8")
        except UnicodeDecodeError:
            return data.decode("latin-1", errors="ignore")

    elif suffix == ".csv":
        try:
            frame = pd.read_csv(io.BytesIO(data))
            # Format dataframe as Markdown table so LLMs understand columns & headers
            return frame.to_markdown(index=False) or ""
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Failed to parse CSV file: {str(e)}",
            )

    elif suffix == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            pages_text = [page.extract_text() or "" for page in reader.pages]
            return "\n\n".join(pages_text).strip()
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Failed to extract text from PDF: {str(e)}",
            )

    elif suffix == ".docx":
        try:
            from docx import Document
            document = Document(io.BytesIO(data))
            paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
            # Also extract table text
            for table in document.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)
            return "\n\n".join(paragraphs)
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Failed to extract text from DOCX: {str(e)}",
            )

    elif suffix == ".pptx":
        try:
            from pptx import Presentation
            presentation = Presentation(io.BytesIO(data))
            lines: list[str] = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_lines: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_lines.append(shape.text.strip())
                if slide_lines:
                    lines.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_lines))
            return "\n\n".join(lines)
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
                detail=f"Failed to extract text from PPTX: {str(e)}",
            )

    return ""


def chunk_text(text: str, *, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    """
    Sliding Window Chunking Algorithm.
    Splits text into windows of `chunk_size` characters with `overlap` character overlap.
    """
    cleaned = " ".join(text.split())
    if not cleaned:
        return []

    chunks: list[str] = []
    start = 0
    text_len = len(cleaned)

    while start < text_len:
        end = min(start + chunk_size, text_len)
        chunk = cleaned[start:end].strip()
        if chunk:
            chunks.append(chunk)

        if end == text_len:
            break

        # Move forward by (chunk_size - overlap)
        start = max(end - overlap, start + 1)

    return chunks
